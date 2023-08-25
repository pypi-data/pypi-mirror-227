from pptx import Presentation
from pptx.parts.chart import ChartPart
from pptx.chart.data import CategoryChartData
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.table import _Cell,Table, _Row, _Column
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

from cast_common.logger import Logger, INFO,DEBUG
from os.path import abspath,exists
from copy import deepcopy

from pandas import DataFrame

class PowerPoint():

    log = None

    def __init__(self,template:str,output:str,log_level=INFO) -> None:

        if PowerPoint.log is None:
            PowerPoint.log = Logger('PowerPoint',log_level)

        self._template_name = abspath(template)
        if not exists(self._template_name):
            raise ValueError(f'Template {self._template_name} not found!')

        self._output = abspath(output)
        self.log.info(f'Generating deck {self._output} from {self._template_name}')

        self._prs = Presentation(self._template_name)

    def save(self):
        self._prs.save(self._output)

    def delete_run(self,run):
        r = run._r
        r.getparent().remove(r)

    def get_shape_by_name(self, name, use_slide=None):
        slides = self._prs.slides
        if use_slide != None:
            slides = [use_slide] 

        for slide in slides:
            for shape in slide.shapes:
                parts = shape.name.split(':')
                shape_name = parts[0]
                if shape_name == name:
                    return shape
        return None

    """ **************************************************************************************************************
                                            Text Replace Functionality 
    ************************************************************************************************************** """
    def replace_text (self, search_str:str, repl:any, tbd_for_blanks=True,slide=None):
        if tbd_for_blanks:
            skip = False
            omit_list = ["immediate_action","other","risk_detail"]
            for s in omit_list:
                if s in search_str:
                    skip=True

            if repl is not None and (type(repl)==int or type(repl)==float):
                repl=f'{repl:,}'

            if not skip and len(str(repl)) == 0:
                repl = 'TBD'

        if slide is None:
            for s in self._prs.slides:
                self._replace_slide_text(s, search_str, repl)
        else:
            self._replace_slide_text(slide, search_str, repl)

    def _replace_slide_text (self, slide, search_str, repl):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if search_str in paragraph.text:
                        self._replace_paragraph_text(paragraph,search_str,repl)
            elif shape.has_table:
                tbl=shape.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0,row_count):
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        for paragraph in cell.text_frame.paragraphs:
                            self._replace_paragraph_text(paragraph,search_str,repl)

    def _replace_paragraph_text (self, paragraph, search_str, repl):
        """
            search all the runs in a paragraph and replace the search_str with repl
        """
        if search_str in paragraph.text:
            t_parags = len(paragraph.runs)
            for run_idx in range(t_parags):
                run = paragraph.runs[run_idx]
                if run.text.count('{')==run.text.count('}') and search_str in run.text:
                    run.text = run.text.replace(str(search_str), str(repl))
                elif run.text.count('{')!=run.text.count('}'):
                    #have a partial tag, need to merge runs
                    base_run = run
                    close_found = False
                    if run_idx < t_parags:
                        for mrg_idx in range(run_idx+1,t_parags):
                            m_run = paragraph.runs[mrg_idx]
                            base_run.text = base_run.text + m_run.text
                            if '}' in m_run.text:
                                close_found = True
                                break
                        if close_found:
                            #delete all extra runs
                            for i in reversed(range(run_idx+1,mrg_idx+1)):
                               self.delete_run(paragraph.runs[i]) 
                            self._replace_paragraph_text(paragraph, search_str, repl)
                            break

    def _merge_runs(self, paragraph):
        cur_text=''
        first=True
        for run in paragraph.runs:
            cur_text = cur_text + run.text
            if first != True:
                self.delete_run(run)
            first=False
        if len(paragraph.runs) > 0:   
            run = paragraph.runs[0]
            run.text = cur_text
        else: 
            run = paragraph.add_run()
        return run

    def update_chart(self, name:str,df:DataFrame):
        shape = self.get_shape_by_name(name)
        if shape != None:
            titles = list(df.index.values)
            data=df[df.columns[0]].values.tolist()
            for i in range(0,len(data)):
                if (isinstance(data[i],str)):
                    data[i] = int(data[i].replace(',',''))

            if shape.has_chart:
                chart_data = CategoryChartData()
                chart_data.categories = titles
                chart_data.add_series('Series 1', data)
                shape.chart.replace_data(chart_data)

    def replace_block(self, begin_tag, end_tag, repl_text):
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text
                        if text.find(begin_tag)!=-1:
                            run=self.merge_runs(paragraph)
                            run_text = run.text
                            text_prefix = text[:run_text.find(begin_tag)]
                            text_suffix = text[run_text.find(end_tag)+len(end_tag):]
                            new_text = text_prefix + repl_text + text_suffix
                            run.text = run.text.replace(run_text,new_text)


    """ ***********************************************************************************************************
                                                Update Table Functionality
    *********************************************************************************************************** """
    def update_table(self, name, df:DataFrame, include_index=True, 
                     background=None, forground=None,neg_num_bkg:list=None, 
                     has_header=True,max_rows=-1):
        table_shape = self.get_shape_by_name(name)
        
        #a notation of :## at the end of the table name indicates the max row count
        splt = table_shape.name.split(':')
        if len(splt) > 1:
            max_rows=int(splt[1])
            if max_rows == 0:
                max_rows=-1
                self.log.warning(f'Invalid max rows notation: {table_shape.name}')

        if table_shape != None and table_shape.has_table:
            table=table_shape.table

            colnames = list(df.columns)
            self.log.debug(f'filling table {name} with {len(df.index)} rows of data')

            # are there enough rows 
            rows, cols = df.shape

            trows = len(table._tbl.tr_lst)
            if not has_header:
                trows=trows+1
            drows = len(df.index)
            if max_rows > 0:
                drows = min(rows,max_rows)
            
            if trows-1 < drows:
                nrc = drows-trows+1
                for r in range(nrc):
                    self.add_row(table)


            # Insert the row zero names
            if include_index:
                for col_index, col_name in enumerate(df.index):
                    try:
                        if has_header:
                            cell = table.cell(col_index+1,0)
                        else:
                            cell = table.cell(col_index,0)
                        text = str(col_name)
                        run = self._merge_runs(cell.text_frame.paragraphs[0])
                        run.text = text
                    except IndexError:
                        self.log.debug(f'index error in update_table ({name} @ row {col_index} with {text}) while setting df index')

            if background:
                cols = cols-1
            if forground:
                cols = cols-1

            if background:
                self.set_table_bgcolor(table,df[background],rows,cols,has_header)
            if forground:
                try:
                    self.set_table_font_color(table,df[forground],rows,cols,has_header)
                except (KeyError):
                    self.log.warning(f'error setting forground for {name}')

            m = df.values
            for row in range(rows):
                for col in range(cols):
                    val = m[row, col]
                    text = str(val)
                    if include_index:
                        tbl_col=col+1
                    else:
                        tbl_col=col

                    try:
                        if has_header:
                            cell = table.cell(row+1,tbl_col)
                        else:
                            cell = table.cell(row,tbl_col)

                        run = self._merge_runs(cell.text_frame.paragraphs[0])
                        run.text = text

                        test = 0
                        if type(val)==str and PowerPoint.isFloat(val.replace('%','')):
                            test = float(val.replace('%','').strip())

                        if (type(val)==str and PowerPoint.isFloat(test) and test < 0) or \
                           ((type(val)==int or type(val)==float) and int(val) < 0):
                            text_frame = cell.text_frame
                            run = self._merge_runs(text_frame.paragraphs[0])
                            run.font.color.rgb=RGBColor(255,255,255)
                            p = text_frame.paragraphs[0]

                            rPr = run._r.get_or_add_rPr()
                            hl = OxmlElement("a:highlight")
                            srgbClr = OxmlElement("a:srgbClr")

                            setattr(srgbClr,'val','FF0000')
                            hl.append(srgbClr)
                            rPr.append(hl)


                            pass

                    except IndexError:
                        self.log.debug(f'index error in update_table ({name}) while setting values')
    
    def set_table_bgcolor(self,table,colors,rows,cols,has_header):
        for row in range(rows):
            rgb = colors.iloc[row].split(",")
            for col in range(cols):
                try:
                    if has_header:
                        cell = table.cell(row+1,col)
                    else:
                        cell = table.cell(row,col)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(int(rgb[0]), int(rgb [1]), int(rgb[2]))
                except IndexError:
                    self.warning('index error in update_table while setting background color')
    
    def set_table_font_color(self,table,colors,rows,cols,has_header):
        for row in range(rows):
            rgb = colors.iloc[row].split(",")
            for col in range(cols):
                try:
                    if has_header:
                        cell = table.cell(row+1,col)
                    else:
                        cell = table.cell(row,col)

                    paragraph = cell.text_frame.paragraphs[0]
                    run = self._merge_runs(paragraph)
                    run.font.color.rgb=RGBColor(int(rgb[0]), int(rgb [1]), int(rgb[2]))
                except IndexError:
                    self.log.warning('index error in update_table while setting background color')

    def change_paragraph_color(self,paragraph,rgb):
        run = self._merge_runs(paragraph)
        run.font.color.rgb=RGBColor(int(rgb[0]), int(rgb [1]), int(rgb[2]))

    def isFloat(str):
        rslt = True
        try:
            float(str)
        except ValueError:
            rslt = False
        return rslt 

    def add_row(self,table: Table) -> _Row:
        new_row = deepcopy(table._tbl.tr_lst[-1]) 
        # duplicating last row of the table as a new row to be added

        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            run = self._merge_runs(cell.text_frame.paragraphs[0])
            run.text = 'xxx' # defaulting cell contents to empty text

        table._tbl.append(new_row) 
        return new_row


    def duplicate_slides(self, app_cnt):
        for cnt in range(2,app_cnt+1):
            for idx, slide in enumerate(self._prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text == "{app_per_page}":
                                new_slide = self.copy_slide(idx)
                                self._replace_slide_text(new_slide,"{app_per_page}","")
                                self._replace_slide_text(new_slide,"{app1_",f'{{app{cnt}_')
                                self.replace_shape_name(new_slide,"app1_",f'app{cnt}_')
                            # if paragraph.text == "{multi_page}":
                            #     new_slide = self.copy_slide(idx)
                            #     self.replace_slide_text(new_slide,"{app_per_page}","")
                            #     self.replace_slide_text(new_slide,"{app1_",f'{{app{cnt}_')
                            #     self.replace_shape_name(new_slide,"app1_",f'app{cnt}_')

        for idx, slide in enumerate(self._prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text == "{app_per_page}":
                            self._replace_slide_text(slide,"{app_per_page}","")

    def copy_slide(self,index=-1,template=None):
        if index<0 and template is None:
            raise KeyError('invalid parameters: either index or template are required')

        if template is not None:
            source = template
        else:
            source = self._prs.slides[index]

        blank_slide_layout = source.slide_layout
        dest = self._prs.slides.add_slide(blank_slide_layout)

        for shp in source.shapes:
            el = shp.element
            newel = deepcopy(el)
            dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for key, value in source.part.rels.items():
            try:
                # Make sure we don't copy a notesSlide relation as that won't exist
                if "notesSlide" not in value.reltype:
                    target = value._target
                    # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
                    if "chart" in value.reltype:
                        partname = target.package.next_partname(
                            ChartPart.partname_template)
                        xlsx_blob = target.chart_workbook.xlsx_part.blob
                        target = ChartPart(partname, target.content_type,
                                        deepcopy(target._element), package=target.package)

                        target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(
                            xlsx_blob, target.package)

                    dest.part.rels.add_relationship(value.reltype,
                                                    target,
                                                    value.rId)
            except AttributeError as err:
                self.logger.error(f'Attribute Error {err} while copying slide {index} part {key}')
        return dest

    def copy_block(self, tag, prefix, count,slide=None):
        search_start = f'{{{tag}}}'
        search_end = f'{{end_{tag}}}'

        block = []

        slides = []
        if slide is not None:
            slides = [slide] 
        else:
            slides = self._prs.slides


        found=False
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text=paragraph.text
                        if not found and text.find(search_start)!=-1:
                            #is the end in the same paragraph?
                            if text.find(search_end)!=-1:
                                run=self._merge_runs(paragraph)
                                old_text = run.text
                                text_prefix = text[:old_text.find(search_end)]
                                text_suffix = text[old_text.find(search_end):]

                                sub_text = ""
                                for app_no in range (2,count+1):
                                    temp = old_text[old_text.find(search_start)+len(search_start):old_text.find(search_end)]
                                    for p in prefix:
                                        temp = temp.replace(f'{p}1',f'{p}{app_no}')
                                    sub_text = sub_text + ", " + temp
                                
                                new_text = text_prefix.replace(search_start,'') + sub_text + text_suffix.replace(search_end,'')
                                #new_text = text_prefix + sub_text + text_suffix

                                run.text = run.text.replace(old_text,new_text)
                            else:
                                found=True
                        elif found and text.find(search_end)!=-1:
                            found=False
                            for app_no in range(2,count+1):
                                self.paste_block(block, shape.text_frame,app_no)
                            if paragraph.text==search_end:
                                self.delete_paragraph(paragraph)
                            block=[]
                        if found:
                            if paragraph.text==search_start:
                                self.delete_paragraph(paragraph)
                            else:
                                block.append(paragraph)

    def delete_paragraph(self,paragraph):
        p = paragraph._p
        parent_element = p.getparent()
        parent_element.remove(p)

    def paste_block(self,block, text_frame,app_no):
        start_at = block[-1]

        for b in block:
            p = text_frame.add_paragraph()
            p.alignment = b.alignment
            p.line_spacing = b.line_spacing
            p.level = b.level

            for r in b.runs:    
                run = p.add_run()
                run.text = deepcopy(r.text)
                font = run.font
                font.name = r.font.name
                font.size = r.font.size
                font.bold = r.font.bold
                font.italic = r.font.italic
                if hasattr(r.font.color, 'rgb'):
                    font.color.rgb = r.font.color.rgb
                # else:
                #     font.color.theme_color = r.font.color.theme_color 

            run = self.merge_runs(p)
            run.text = run.text.replace("{app1_",f'{{app{app_no}_')
            run.text = run.text.replace("{end_app1_",f'{{end_app{app_no}_')

    def merge_runs(self, paragraph):
        cur_text=''
        first=True
        for run in paragraph.runs:
            cur_text = cur_text + run.text
            if first != True:
                self.delete_run(run)
            first=False
        if len(paragraph.runs) > 0:   
            run = paragraph.runs[0]
            run.text = cur_text
        else: 
            run = paragraph.add_run()
        return run

    def replace_shape_name (self, slide, search_str, repl_str):
        for shape in slide.shapes:
            if shape.name.find(search_str) != -1: 
                shape.name = shape.name.replace(search_str,repl_str)

    def get_page_no(self,shape):
        page_no = 0
        if shape:
            while True:
                shape = self.get_shape_parent(shape)
                if type(shape).__name__ == 'Slide':
                    break;
            
            page_no = self._prs.slides.index(shape) 
        return page_no

    def get_shape_parent(self,shape):
        rslt = None
        if hasattr(shape,'_parent'):
            rslt = shape._parent
        return rslt

